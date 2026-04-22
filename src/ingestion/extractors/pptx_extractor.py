"""
pptx_extractor.py — Extractor for .pptx files.

Extracts text from each slide as a single block: title, body, and
speaker notes. One block per slide preserves the coherence of each slide.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt


def _extract_slide_text(slide) -> str:
    """
    Extract all text from a single slide including notes.

    Args:
        slide: A python-pptx Slide object.

    Returns:
        Text content of the slide as a single string.
    """
    blocks = []

    # extract title separately so it leads the block
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            blocks.append(f"Title: {title_text}")

    # extract body text from all other shapes
    for shape in slide.shapes:
        if shape == slide.shapes.title:        # already handled above
            continue
        if not shape.has_text_frame:           # skip images, charts, etc.
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                blocks.append(text)

    # extract speaker notes — often contain useful context
    if slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame:
            notes_text = notes_frame.text.strip()
            if notes_text:
                blocks.append(f"Notes: {notes_text}")

    return "\n".join(blocks)


def extract(file_path: Path) -> str:
    """
    Extract plain text from a PowerPoint file.

    Args:
        file_path: Path to the .pptx file.

    Returns:
        Extracted text with one block per slide, separated by blank lines.

    Raises:
        ValueError: If the file extension is not .pptx.
    """
    if file_path.suffix.lower() != ".pptx":
        raise ValueError(f"pptx_extractor only handles .pptx, got: {file_path.suffix}")

    presentation = Presentation(file_path)
    slides = []

    for i, slide in enumerate(presentation.slides, start=1):
        slide_text = _extract_slide_text(slide)
        if slide_text.strip():                 # skip completely empty slides
            slides.append(f"--- Slide {i} ---\n{slide_text}")

    return "\n\n".join(slides)


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python -m src.ingestion.extractors.pptx_extractor <file.pptx>")
        sys.exit(1)

    target = Path(sys.argv[1])

    if not target.exists():
        print(f"Error: file not found — {target}")
        sys.exit(1)

    result = extract(target)

    print(f"--- Extracted {len(result):,} characters from {target.name} ---\n")
    print(result[:500])